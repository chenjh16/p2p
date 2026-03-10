"""End-to-end tests for --continue-run functionality."""

from __future__ import annotations

import json
import os

import pytest
from pptx import Presentation

from tests.conftest import MOCK_BASE_URL, create_partial_run


def test_continue_run_postprocess_only(three_page_pdf, tmp_path, monkeypatch):
    """--continue-run with 'post-process' choice should assemble PPTX from existing slides."""
    os.chdir(tmp_path)
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    run_dir = create_partial_run(tmp_path, three_page_pdf, pages_received=[0, 1], total_pages=3)

    inputs = iter(["p"])
    monkeypatch.setattr("builtins.input", lambda _: next(inputs))

    from src.continue_run import run_continue

    run_continue(run_dir)

    output_files = [f for f in os.listdir(run_dir) if f.endswith(".pptx")]
    assert len(output_files) == 1

    pptx_path = os.path.join(run_dir, output_files[0])
    prs = Presentation(pptx_path)
    assert len(prs.slides) == 2


def test_continue_run_all_present(three_page_pdf, tmp_path, monkeypatch):
    """--continue-run when all slides present should assemble PPTX directly."""
    os.chdir(tmp_path)
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    run_dir = create_partial_run(tmp_path, three_page_pdf, pages_received=[0, 1, 2], total_pages=3)

    inputs = iter(["y"])
    monkeypatch.setattr("builtins.input", lambda _: next(inputs))

    from src.continue_run import run_continue

    run_continue(run_dir)

    output_files = [f for f in os.listdir(run_dir) if f.endswith(".pptx")]
    assert len(output_files) == 1

    pptx_path = os.path.join(run_dir, output_files[0])
    prs = Presentation(pptx_path)
    assert len(prs.slides) == 3


def test_continue_run_quit(three_page_pdf, tmp_path, monkeypatch):
    """--continue-run with 'quit' choice should exit without producing PPTX."""
    os.chdir(tmp_path)
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    run_dir = create_partial_run(tmp_path, three_page_pdf, pages_received=[0], total_pages=3)

    inputs = iter(["q"])
    monkeypatch.setattr("builtins.input", lambda _: next(inputs))

    from src.continue_run import run_continue

    run_continue(run_dir)

    output_files = [f for f in os.listdir(run_dir) if f.endswith(".pptx")]
    assert len(output_files) == 0


def test_continue_run_generate_missing(mock_server, three_page_pdf, tmp_path, monkeypatch):
    """--continue-run with 'continue' should generate missing pages via API and assemble."""
    os.chdir(tmp_path)
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    run_dir = create_partial_run(tmp_path, three_page_pdf, pages_received=[0], total_pages=3)

    meta_path = os.path.join(run_dir, "metadata.json")
    with open(meta_path) as f:
        meta = json.load(f)
    meta["runtime_params"]["api_provider"] = "openai"
    meta["runtime_params"]["skip_postprocess"] = True
    with open(meta_path, "w") as f:
        json.dump(meta, f)

    inputs = iter(["c"])
    monkeypatch.setattr("builtins.input", lambda _: next(inputs))

    monkeypatch.setenv("OPENAI_BASE_URL", MOCK_BASE_URL)
    monkeypatch.setenv("OPENAI_API_KEY", "test-key")

    from src.continue_run import run_continue

    run_continue(run_dir)

    slides_dir = os.path.join(run_dir, "slides")
    slide_files = [f for f in os.listdir(slides_dir) if f.endswith(".xml")]
    assert len(slide_files) >= 2

    output_files = [f for f in os.listdir(run_dir) if f.endswith(".pptx")]
    assert len(output_files) == 1


def test_continue_run_missing_dir(tmp_path, monkeypatch):
    """--continue-run with non-existent directory should exit with error."""
    os.chdir(tmp_path)
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    with pytest.raises(SystemExit):
        from src.continue_run import run_continue

        run_continue(str(tmp_path / "nonexistent-run-dir"))
