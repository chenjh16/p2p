"""Unit tests for individual modules."""

from __future__ import annotations

import json
import os
import shutil

import fitz
import pytest

# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def sample_pdf(tmp_path):
    """Create a minimal 2-page PDF."""
    pdf_path = str(tmp_path / "test.pdf")
    doc = fitz.open()
    for i in range(2):
        page = doc.new_page(width=720, height=405)
        page.insert_text((100, 200), f"Page {i}", fontsize=24)
    doc.save(pdf_path)
    doc.close()
    return pdf_path


@pytest.fixture()
def three_page_pdf(tmp_path):
    """Create a 3-page PDF for batch testing."""
    pdf_path = str(tmp_path / "three.pdf")
    doc = fitz.open()
    for i in range(3):
        page = doc.new_page(width=720, height=405)
        page.insert_text((100, 200), f"Slide {i}", fontsize=20)
    doc.save(pdf_path)
    doc.close()
    return pdf_path


# ---------------------------------------------------------------------------
# pdf_preprocessor
# ---------------------------------------------------------------------------


class TestPdfPreprocessor:
    """Tests for src.pdf_preprocessor."""

    def test_renders_correct_page_count(self, sample_pdf):
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        assert len(pages) == 2

    def test_returns_png_bytes(self, sample_pdf):
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        img_bytes, _ = pages[0]
        assert isinstance(img_bytes, bytes)
        assert img_bytes[:8] == b"\x89PNG\r\n\x1a\n"

    def test_metadata_fields(self, sample_pdf):
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        _, meta = pages[0]
        assert meta["page_num"] == 0
        assert meta["width_pt"] == pytest.approx(720.0, abs=1)
        assert meta["height_pt"] == pytest.approx(405.0, abs=1)
        assert "width_px" in meta
        assert "height_px" in meta

    def test_dpi_affects_image_size(self, sample_pdf):
        from src.pdf_preprocessor import pdf_to_images

        pages_low = pdf_to_images(sample_pdf, dpi=72)
        pages_high = pdf_to_images(sample_pdf, dpi=144)
        assert len(pages_low[0][0]) < len(pages_high[0][0])

    def test_page_numbers_sequential(self, three_page_pdf):
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(three_page_pdf, dpi=72)
        nums = [meta["page_num"] for _, meta in pages]
        assert nums == [0, 1, 2]


class TestSnapSlideDimensions:
    """Tests for snap_slide_dimensions aspect ratio detection."""

    def test_detects_16_9(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(720, 405)
        assert label == "16:9"
        assert w == 720
        assert h == 405

    def test_detects_4_3(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(720, 540)
        assert label == "4:3"
        assert w == 720
        assert h == 540

    def test_detects_16_10(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(720, 450)
        assert label == "16:10"
        assert w == 720
        assert h == 450

    def test_snaps_close_ratio(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(722, 406)
        assert label == "16:9"
        assert w == 720
        assert h == 405

    def test_custom_ratio_preserved(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(500, 500)
        assert label == "custom"
        assert w == 500
        assert h == 500

    def test_zero_height(self):
        from src.pdf_preprocessor import snap_slide_dimensions

        w, h, label = snap_slide_dimensions(720, 0)
        assert label == "unknown"


# ---------------------------------------------------------------------------
# message_builder
# ---------------------------------------------------------------------------


class TestMessageBuilder:
    """Tests for src.message_builder."""

    def test_builds_system_and_user_messages(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages, enable_animations=False)

        assert len(messages) == 2
        assert messages[0]["role"] == "system"
        assert messages[1]["role"] == "user"

    def test_system_prompt_contains_key_sections(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages)

        system = messages[0]["content"]
        assert "PresentationML" in system
        assert "write_slide_xml" in system
        assert "Font Size Calibration" in system
        assert "Atomic Reconstruction" in system

    def test_animation_section_included_when_enabled(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)

        msgs_off = build_messages(pages, enable_animations=False)
        msgs_on = build_messages(pages, enable_animations=True)

        assert "Morph" not in msgs_off[0]["content"]
        assert "Morph" in msgs_on[0]["content"]

    def test_user_message_contains_images(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages)

        user_content = messages[1]["content"]
        image_parts = [p for p in user_content if p.get("type") == "image_url"]
        assert len(image_parts) == 2

    def test_task_instruction_at_end(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages)

        user_content = messages[1]["content"]
        last_text = user_content[-1]
        assert last_text["type"] == "text"
        assert "CRITICAL" in last_text["text"]
        assert "exactly 2 parallel tool calls" in last_text["text"]
        assert "Atomic Reconstruction" in last_text["text"]
        assert "font size reduction" in last_text["text"]

    def test_chinese_prompt(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages, prompt_lang="zh")

        system = messages[0]["content"]
        assert "演示文稿重建引擎" in system
        assert "原子化重建" in system

    def test_slide_dimensions_in_message(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages)

        user_content = messages[1]["content"]
        text_parts = [p["text"] for p in user_content if p.get("type") == "text"]
        combined = " ".join(text_parts)
        assert "720" in combined
        assert "405" in combined


# ---------------------------------------------------------------------------
# token_estimator
# ---------------------------------------------------------------------------


class TestTokenEstimator:
    """Tests for src.token_estimator."""

    def test_estimates_text_tokens(self):
        from src.token_estimator import estimate_tokens

        messages = [{"role": "system", "content": "You are a helpful assistant."}]
        result = estimate_tokens(messages, model="gpt-4")
        assert result["text_tokens"] > 0
        assert result["image_count"] == 0
        assert result["image_tokens"] == 0

    def test_counts_images(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images
        from src.token_estimator import estimate_tokens

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages)
        result = estimate_tokens(messages)

        assert result["image_count"] == 2
        assert result["image_tokens"] > 0
        assert result["total_input_tokens"] == result["text_tokens"] + result["image_tokens"]

    def test_cost_estimate_structure(self):
        from src.token_estimator import estimate_tokens

        messages = [{"role": "user", "content": "Hello"}]
        result = estimate_tokens(messages)

        cost = result["estimated_cost_usd"]
        assert isinstance(cost, dict)
        assert "input_cost_usd" in cost
        assert "output_cost_usd" in cost
        assert "total_cost_usd" in cost
        assert cost["total_cost_usd"] >= 0

    def test_output_estimate_scales_with_images(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images
        from src.token_estimator import estimate_tokens

        pages = pdf_to_images(sample_pdf, dpi=72)
        result_2 = estimate_tokens(build_messages(pages))

        result_1 = estimate_tokens(build_messages(pages[:1]))

        assert result_2["estimated_output_tokens"] > result_1["estimated_output_tokens"]

    def test_estimated_response_time(self):
        from src.token_estimator import estimate_tokens

        messages = [{"role": "user", "content": "Hello"}]
        result = estimate_tokens(messages, reasoning_effort="medium")

        assert "assumed_output_tps" in result
        assert result["assumed_output_tps"] == 30.0
        assert result["reasoning_effort"] == "medium"
        assert result["reasoning_multiplier"] == 1.5
        assert "estimated_response_time_seconds" in result
        expected_time = (result["estimated_output_tokens"] / 30.0) * 1.5
        assert abs(result["estimated_response_time_seconds"] - round(expected_time, 1)) < 0.2

    def test_estimated_response_time_high_reasoning(self):
        from src.token_estimator import estimate_tokens

        messages = [{"role": "user", "content": "Hello"}]
        result_high = estimate_tokens(messages, reasoning_effort="high")
        result_low = estimate_tokens(messages, reasoning_effort="low")

        assert result_high["reasoning_multiplier"] == 2.5
        assert result_low["reasoning_multiplier"] == 1.0
        assert result_high["estimated_response_time_seconds"] > result_low["estimated_response_time_seconds"]

    def test_openai_image_tokens_high_detail(self):
        from src.token_estimator import _openai_image_tokens

        # 2880×1620 at high detail: scaled to 768 shortest side → 1365×768, tiles 3×2=6
        tokens = _openai_image_tokens(2880, 1620, "high")
        assert tokens == 85 + 170 * 6  # 1105

    def test_openai_image_tokens_low_detail(self):
        from src.token_estimator import _openai_image_tokens

        tokens = _openai_image_tokens(2880, 1620, "low")
        assert tokens == 85

    def test_anthropic_image_tokens(self):
        from src.token_estimator import _anthropic_image_tokens

        # 2880×1620: long edge 2880 > 1568 → scale to 1568×882, then area/750
        tokens = _anthropic_image_tokens(2880, 1620)
        assert tokens > 0
        assert tokens < 2000  # should be ~1843

    def test_anthropic_image_tokens_small(self):
        from src.token_estimator import _anthropic_image_tokens

        # 200×200: no resize needed, tokens = 200*200/750 ≈ 54
        tokens = _anthropic_image_tokens(200, 200)
        assert tokens == 54

    def test_anthropic_vs_openai_different_counts(self):
        from src.token_estimator import _anthropic_image_tokens, _openai_image_tokens

        openai_tokens = _openai_image_tokens(1920, 1080, "high")
        anthropic_tokens = _anthropic_image_tokens(1920, 1080)
        assert openai_tokens != anthropic_tokens


# ---------------------------------------------------------------------------
# xml_validator
# ---------------------------------------------------------------------------


class TestXmlValidator:
    """Tests for src.xml_validator."""

    GOOD_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
    <p:grpSpPr/>
  </p:spTree></p:cSld>
</p:sld>"""

    def test_valid_xml_passes_through(self):
        from src.xml_validator import validate_and_fix

        result = validate_and_fix(self.GOOD_XML, 0)
        assert "<p:sld" in result

    def test_strips_markdown_fences(self):
        from src.xml_validator import validate_and_fix

        fenced = f"```xml\n{self.GOOD_XML}\n```"
        result = validate_and_fix(fenced, 0)
        assert "```" not in result
        assert "<p:sld" in result

    def test_adds_xml_declaration(self):
        from src.xml_validator import validate_and_fix

        no_decl = self.GOOD_XML.replace('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n', "")
        result = validate_and_fix(no_decl, 0)
        assert result.startswith("<?xml")

    def test_fixes_unescaped_ampersand(self):
        from src.xml_validator import validate_and_fix

        xml_with_amp = self.GOOD_XML.replace('name=""', 'name="A &amp; B"')
        result = validate_and_fix(xml_with_amp, 0)
        assert "<p:sld" in result

    def test_broken_xml_returns_fallback(self):
        from src.xml_validator import validate_and_fix

        result = validate_and_fix("<p:sld><broken", 5)
        assert "ErrorInfo" in result
        assert "page 5" in result

    def test_adds_missing_namespace(self):
        from src.xml_validator import validate_and_fix

        xml_no_ns = '<p:sld><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld></p:sld>'
        result = validate_and_fix(xml_no_ns, 0)
        assert "<p:sld" in result


# ---------------------------------------------------------------------------
# system_prompt
# ---------------------------------------------------------------------------


class TestSystemPrompt:
    """Tests for src.system_prompt."""

    def test_get_system_prompt_en(self):
        from src.system_prompt import get_system_prompt

        prompt = get_system_prompt("en")
        assert "presentation reconstruction engine" in prompt
        assert "Atomic Reconstruction" in prompt

    def test_get_system_prompt_zh(self):
        from src.system_prompt import get_system_prompt

        prompt = get_system_prompt("zh")
        assert "演示文稿重建引擎" in prompt
        assert "原子化重建" in prompt

    def test_get_animation_section_en(self):
        from src.system_prompt import get_animation_section

        section = get_animation_section("en")
        assert "Morph Transition" in section

    def test_get_animation_section_zh(self):
        from src.system_prompt import get_animation_section

        section = get_animation_section("zh")
        assert "Morph 转场" in section

    def test_default_is_english(self):
        from src.system_prompt import SYSTEM_PROMPT, SYSTEM_PROMPT_EN

        assert SYSTEM_PROMPT == SYSTEM_PROMPT_EN

    def test_tool_definition_structure(self):
        from src.system_prompt import WRITE_SLIDE_XML_TOOL

        assert WRITE_SLIDE_XML_TOOL["type"] == "function"
        func = WRITE_SLIDE_XML_TOOL["function"]
        assert func["name"] == "write_slide_xml"
        params = func["parameters"]
        assert "page_num" in params["properties"]
        assert "slide_xml" in params["properties"]
        assert params["required"] == ["page_num", "slide_xml"]

    def test_anthropic_tool_definition_structure(self):
        from src.system_prompt import WRITE_SLIDE_XML_TOOL_ANTHROPIC

        assert "name" in WRITE_SLIDE_XML_TOOL_ANTHROPIC
        assert WRITE_SLIDE_XML_TOOL_ANTHROPIC["name"] == "write_slide_xml"
        assert "description" in WRITE_SLIDE_XML_TOOL_ANTHROPIC
        assert "input_schema" in WRITE_SLIDE_XML_TOOL_ANTHROPIC
        schema = WRITE_SLIDE_XML_TOOL_ANTHROPIC["input_schema"]
        assert "properties" in schema
        assert "page_num" in schema["properties"]
        assert "slide_xml" in schema["properties"]
        assert "required" in schema
        assert schema["required"] == ["page_num", "slide_xml"]

    def test_font_calibration_in_prompt(self):
        from src.system_prompt import get_system_prompt

        prompt = get_system_prompt("en")
        assert "REDUCE it by 20%" in prompt
        assert "sz=1600" in prompt

    def test_table_requirement_in_prompt(self):
        from src.system_prompt import get_system_prompt

        prompt = get_system_prompt("en")
        assert "NEVER use a raster placeholder for any table" in prompt


# ---------------------------------------------------------------------------
# artifacts
# ---------------------------------------------------------------------------


class TestArtifactStore:
    """Tests for src.artifacts.ArtifactStore."""

    def test_creates_directory_structure(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        assert os.path.isdir(store.root)
        assert os.path.isdir(store.pages_dir)
        assert os.path.isdir(store.slides_dir)
        assert store.root.startswith("runs/run-")
        shutil.rmtree("runs")

    def test_dry_run_prefix(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf, dry_run=True)
        assert store.root.startswith("runs/dry-run-")
        shutil.rmtree("runs")

    def test_replay_prefix(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf, replay_of="runs/run-test-20260101-000000")
        assert store.root.startswith("runs/replay-")
        shutil.rmtree("runs")

    def test_save_page_images(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore
        from src.pdf_preprocessor import pdf_to_images

        store = ArtifactStore(sample_pdf)
        pages = pdf_to_images(sample_pdf, dpi=72)
        store.save_page_images(pages)

        assert os.path.isfile(os.path.join(store.pages_dir, "page_000.png"))
        assert os.path.isfile(os.path.join(store.pages_dir, "page_001.png"))
        shutil.rmtree("runs")

    def test_save_run_params(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_run_params({"dpi": 288, "model": "gpt-5.4"})

        path = os.path.join(store.root, "run_params.json")
        assert os.path.isfile(path)
        with open(path) as f:
            data = json.load(f)
        assert data["dpi"] == 288
        shutil.rmtree("runs")

    def test_save_reasoning(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_reasoning("This is the model's thinking process.", batch_idx=0)

        path = os.path.join(store.root, "reasoning.txt")
        assert os.path.isfile(path)
        with open(path) as f:
            assert "thinking process" in f.read()
        shutil.rmtree("runs")

    def test_save_reasoning_skips_empty(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_reasoning("", batch_idx=0)

        assert not os.path.isfile(os.path.join(store.root, "reasoning.txt"))
        shutil.rmtree("runs")

    def test_save_content_text(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_content_text("Some content output", batch_idx=0)

        path = os.path.join(store.root, "content.txt")
        assert os.path.isfile(path)
        shutil.rmtree("runs")

    def test_save_slide_xmls(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_slide_xmls({0: "<p:sld>slide0</p:sld>", 1: "<p:sld>slide1</p:sld>"})

        assert os.path.isfile(os.path.join(store.slides_dir, "slide_000.xml"))
        assert os.path.isfile(os.path.join(store.slides_dir, "slide_001.xml"))
        shutil.rmtree("runs")

    def test_api_response_indexing(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_api_response({"batch": 0})
        store.save_api_response({"batch": 1})

        assert os.path.isfile(os.path.join(store.root, "api_response.json"))
        assert os.path.isfile(os.path.join(store.root, "api_response_1.json"))
        shutil.rmtree("runs")

    def test_save_metadata(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.save_metadata({"pages": 2, "dpi": 288})

        path = os.path.join(store.root, "metadata.json")
        assert os.path.isfile(path)
        with open(path) as f:
            data = json.load(f)
        assert data["pages"] == 2
        shutil.rmtree("runs")

    def test_copy_input(self, sample_pdf, tmp_path):
        os.chdir(tmp_path)
        from src.artifacts import ArtifactStore

        store = ArtifactStore(sample_pdf)
        store.copy_input(sample_pdf)

        copied = os.path.join(store.root, os.path.basename(sample_pdf))
        assert os.path.isfile(copied)
        shutil.rmtree("runs")


# ---------------------------------------------------------------------------
# pptx_assembler
# ---------------------------------------------------------------------------


class TestPPTXAssembler:
    """Tests for src.pptx_assembler.PPTXAssembler."""

    SLIDE_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
       xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
    <p:grpSpPr/>
    <p:sp>
      <p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="914400" y="914400"/><a:ext cx="7315200" cy="914400"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      </p:spPr>
      <p:txBody>
        <a:bodyPr wrap="square" rtlCol="0"/><a:lstStyle/>
        <a:p><a:r><a:rPr lang="en-US" sz="2400"/><a:t>Hello World</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
  </p:spTree></p:cSld>
</p:sld>"""

    def test_assemble_single_slide(self, tmp_path):
        from pptx import Presentation

        from src.pptx_assembler import PPTXAssembler

        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: self.SLIDE_XML})

        output = str(tmp_path / "test.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 1
        texts = [s.text for s in prs.slides[0].shapes if hasattr(s, "text")]
        assert any("Hello World" in t for t in texts)

    def test_assemble_multiple_slides(self, tmp_path):
        from pptx import Presentation

        from src.pptx_assembler import PPTXAssembler

        slide_xmls = {
            0: self.SLIDE_XML,
            1: self.SLIDE_XML.replace("Hello World", "Second Slide"),
        }
        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble(slide_xmls)

        output = str(tmp_path / "multi.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2

    def test_slide_dimensions(self, tmp_path):
        from pptx import Presentation

        from src.pptx_assembler import PPTXAssembler

        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: self.SLIDE_XML})

        output = str(tmp_path / "dims.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert prs.slide_width / 12700 == pytest.approx(720, abs=1)
        assert prs.slide_height / 12700 == pytest.approx(405, abs=1)

    def test_handles_hyperlinks(self, tmp_path):
        from pptx import Presentation

        from src.pptx_assembler import PPTXAssembler

        xml_with_link = self.SLIDE_XML.replace(
            '<a:rPr lang="en-US" sz="2400"/>',
            '<a:rPr lang="en-US" sz="2400"><a:hlinkClick r:id="rId2"/></a:rPr>',
        )
        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: xml_with_link})

        output = str(tmp_path / "link.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# message_builder (Anthropic)
# ---------------------------------------------------------------------------


class TestMessageBuilderAnthropic:
    """Tests for Anthropic-specific message building."""

    def test_anthropic_messages_no_system_role(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages, provider="anthropic")

        for msg in messages:
            assert msg["role"] != "system", "Anthropic uses a separate system parameter, not system role in messages"

    def test_anthropic_messages_use_image_blocks(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages, provider="anthropic")

        user_content = messages[0]["content"]
        image_blocks = [p for p in user_content if p.get("type") == "image"]
        assert len(image_blocks) == 2
        for block in image_blocks:
            assert "source" in block
            assert block["source"].get("type") == "base64"

    def test_anthropic_task_instruction_present(self, sample_pdf):
        from src.message_builder import build_messages
        from src.pdf_preprocessor import pdf_to_images

        pages = pdf_to_images(sample_pdf, dpi=72)
        messages = build_messages(pages, provider="anthropic")

        user_content = messages[0]["content"]
        text_blocks = [p for p in user_content if p.get("type") == "text"]
        last_text = text_blocks[-1]
        assert "CRITICAL" in last_text["text"]
        assert "parallel tool calls" in last_text["text"]

    def test_get_system_prompt_text(self):
        from src.message_builder import get_system_prompt_text

        prompt = get_system_prompt_text(enable_animations=False, prompt_lang="en")
        assert "presentation reconstruction engine" in prompt
        assert "write_slide_xml" in prompt

    def test_get_system_prompt_text_with_animations(self):
        from src.message_builder import get_system_prompt_text

        prompt = get_system_prompt_text(enable_animations=True)
        assert "Morph" in prompt


# ---------------------------------------------------------------------------
# api_client_anthropic (thinking budget)
# ---------------------------------------------------------------------------


class TestAnthropicThinkingBudget:
    """Tests for _thinking_budget mapping in api_client_anthropic."""

    def test_thinking_budget_mapping(self):
        from src.api_client_anthropic import _thinking_budget

        max_tokens = 128000
        low = _thinking_budget("low", max_tokens)
        medium = _thinking_budget("medium", max_tokens)
        high = _thinking_budget("high", max_tokens)
        xhigh = _thinking_budget("xhigh", max_tokens)
        none_val = _thinking_budget("none", max_tokens)
        empty_val = _thinking_budget("", max_tokens)

        assert low == 0
        assert medium > 0
        assert high > medium
        assert xhigh > high
        assert xhigh == min(128000, max_tokens)
        assert none_val == 0
        assert empty_val == 0


# ---------------------------------------------------------------------------
# logging_config
# ---------------------------------------------------------------------------


class TestLoggingConfig:
    """Tests for src.logging_config."""

    def test_setup_logging_does_not_raise(self):
        from src.logging_config import setup_logging

        setup_logging("WARNING")

    def test_get_logger_returns_logger(self):
        import logging

        from src.logging_config import get_logger

        logger = get_logger("test_module")
        assert isinstance(logger, logging.Logger)
        assert logger.name == "test_module"
