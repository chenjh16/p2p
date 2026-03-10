"""End-to-end tests for the core conversion pipeline with mock LLM servers."""

from __future__ import annotations

import json
import os
import shutil

from pptx import Presentation

from src import ApiConfig
from tests.conftest import MOCK_ANTHROPIC_URL, MOCK_BASE_URL


def test_dry_run(sample_pdf, tmp_path):
    """Dry-run should produce artifacts under runs/ without calling the API."""
    os.chdir(tmp_path)
    from src.dry_run import run_dry
    from src.logging_config import setup_logging

    setup_logging("WARNING")
    output_dir = run_dry(
        pdf_path=sample_pdf,
        dpi=96,
        enable_animations=False,
        model_name="gpt-5.4",
        batch_size=25,
    )

    assert output_dir.startswith("runs/")
    assert os.path.isdir(output_dir)
    assert os.path.isfile(os.path.join(output_dir, "metadata.json"))
    assert os.path.isfile(os.path.join(output_dir, "messages_0.json"))
    assert os.path.isfile(os.path.join(output_dir, "system_prompt.md"))
    assert os.path.isfile(os.path.join(output_dir, "token_estimate.json"))
    assert os.path.isfile(os.path.join(output_dir, "run_params.json"))
    assert os.path.isfile(os.path.join(output_dir, "pages", "page_000.png"))
    assert os.path.isfile(os.path.join(output_dir, "pages", "page_001.png"))
    assert os.path.isfile(os.path.join(output_dir, "test.pdf"))

    with open(os.path.join(output_dir, "metadata.json")) as f:
        meta = json.load(f)
    assert meta["pdf_pages"] == 2
    assert meta["runtime_params"]["dpi"] == 96

    shutil.rmtree("runs")


def test_e2e_conversion(mock_server, sample_pdf, tmp_path):
    """Full conversion with mock server should produce a valid PPTX."""
    os.chdir(tmp_path)
    output_pptx = str(tmp_path / "output.pptx")

    from src.api_client import call_llm
    from src.artifacts import ArtifactStore
    from src.logging_config import setup_logging
    from src.message_builder import build_messages
    from src.pdf_preprocessor import pdf_to_images
    from src.pptx_assembler import PPTXAssembler
    from src.token_estimator import estimate_tokens

    setup_logging("WARNING")

    pages = pdf_to_images(sample_pdf, dpi=96)
    assert len(pages) == 2

    store = ArtifactStore(pdf_path=sample_pdf)
    store.save_page_images(pages)

    messages = build_messages(pages, enable_animations=False)
    store.save_messages(messages)

    token_est = estimate_tokens(messages, model="gpt-5.4")
    assert token_est["image_count"] == 2

    stream_log = os.path.join(store.root, "stream.log")
    result = call_llm(
        messages=messages,
        api_cfg=ApiConfig(api_base_url=MOCK_BASE_URL, api_key="test-key", model_name="mock-gpt-5.4"),
        stream_log_path=stream_log,
        reasoning_effort="",
    )

    assert len(result.slide_xmls) == 2
    assert 0 in result.slide_xmls
    assert 1 in result.slide_xmls
    assert result.response_data["finish_reason"] == "stop"
    assert result.response_data["tool_call_count"] == 2
    assert result.reasoning_text == ""

    store.save_slide_xmls(result.slide_xmls)
    store.save_api_response(result.response_data)

    assembler = PPTXAssembler(
        slide_width_pt=pages[0][1]["width_pt"],
        slide_height_pt=pages[0][1]["height_pt"],
    )
    assembler.assemble(result.slide_xmls)
    assembler.save(output_pptx)

    assert os.path.isfile(output_pptx)
    prs = Presentation(output_pptx)
    assert len(prs.slides) == 2

    slide0 = prs.slides[0]
    texts = [shape.text for shape in slide0.shapes if hasattr(shape, "text")]
    assert any("Test Slide Page 0" in t for t in texts)

    slide1 = prs.slides[1]
    texts = [shape.text for shape in slide1.shapes if hasattr(shape, "text")]
    assert any("Test Slide Page 1" in t for t in texts)

    shutil.rmtree("runs")
    os.remove(output_pptx)


def test_e2e_anthropic_conversion(mock_anthropic_server, sample_pdf, tmp_path):
    """Full conversion with mock Anthropic server should produce a valid PPTX."""
    os.chdir(tmp_path)
    output_pptx = str(tmp_path / "output_anthropic.pptx")

    from src.api_client_anthropic import call_anthropic
    from src.artifacts import ArtifactStore
    from src.logging_config import setup_logging
    from src.message_builder import build_messages, get_system_prompt_text
    from src.pdf_preprocessor import pdf_to_images
    from src.pptx_assembler import PPTXAssembler
    from src.token_estimator import estimate_tokens

    setup_logging("WARNING")

    pages = pdf_to_images(sample_pdf, dpi=96)
    assert len(pages) == 2

    store = ArtifactStore(pdf_path=sample_pdf)
    store.save_page_images(pages)

    messages = build_messages(pages, enable_animations=False, provider="anthropic")
    store.save_messages(messages)

    token_est = estimate_tokens(messages, model="claude-opus-4-6", dpi=96)
    assert token_est["image_count"] == 2

    sys_prompt = get_system_prompt_text(enable_animations=False)
    stream_log = os.path.join(store.root, "stream.log")
    result = call_anthropic(
        messages=messages,
        system_prompt=sys_prompt,
        api_cfg=ApiConfig(api_key="test-key", api_base_url=MOCK_ANTHROPIC_URL, model_name="mock-claude"),
        stream_log_path=stream_log,
        reasoning_effort="",
    )

    assert len(result.slide_xmls) == 2
    assert 0 in result.slide_xmls
    assert 1 in result.slide_xmls

    store.save_slide_xmls(result.slide_xmls)
    store.save_api_response(result.response_data)

    assembler = PPTXAssembler(
        slide_width_pt=pages[0][1]["width_pt"],
        slide_height_pt=pages[0][1]["height_pt"],
    )
    assembler.assemble(result.slide_xmls)
    assembler.save(output_pptx)

    assert os.path.isfile(output_pptx)
    prs = Presentation(output_pptx)
    assert len(prs.slides) == 2

    slide0 = prs.slides[0]
    texts = [shape.text for shape in slide0.shapes if hasattr(shape, "text")]
    assert any("Test Slide Page 0" in t for t in texts)

    shutil.rmtree("runs")
    os.remove(output_pptx)


def test_dry_run_anthropic(sample_pdf, tmp_path):
    """Dry-run with Anthropic provider should produce correct artifacts."""
    os.chdir(tmp_path)
    from src.dry_run import run_dry
    from src.logging_config import setup_logging

    setup_logging("WARNING")
    output_dir = run_dry(
        pdf_path=sample_pdf,
        dpi=96,
        enable_animations=False,
        model_name="claude-opus-4-6",
        batch_size=25,
        provider="anthropic",
    )

    assert output_dir.startswith("runs/")
    assert os.path.isdir(output_dir)
    assert os.path.isfile(os.path.join(output_dir, "metadata.json"))
    assert os.path.isfile(os.path.join(output_dir, "messages_0.json"))

    with open(os.path.join(output_dir, "metadata.json")) as f:
        meta = json.load(f)
    assert meta["runtime_params"]["api_provider"] == "anthropic"
    assert meta["runtime_params"]["recommended_batch_size"] > 0
    assert meta["runtime_params"]["gateway_timeout_seconds"] == 600

    with open(os.path.join(output_dir, "messages_0.json")) as f:
        msgs = json.load(f)
    assert all(m.get("role") != "system" for m in msgs)

    shutil.rmtree("runs")


def test_dry_run_custom_output_tps(sample_pdf, tmp_path):
    """Dry-run with custom output_tps should reflect in metadata."""
    os.chdir(tmp_path)
    from src.dry_run import run_dry
    from src.logging_config import setup_logging

    setup_logging("WARNING")
    output_dir = run_dry(
        pdf_path=sample_pdf,
        dpi=96,
        enable_animations=False,
        model_name="gpt-5.4",
        batch_size=25,
        output_tps=100.0,
    )

    with open(os.path.join(output_dir, "metadata.json")) as f:
        meta = json.load(f)
    assert meta["assumed_output_tps"] == 100.0

    with open(os.path.join(output_dir, "token_estimate.json")) as f:
        est = json.load(f)
    assert est["assumed_output_tps"] == 100.0

    shutil.rmtree("runs")


def test_multi_batch_conversion(mock_server, tmp_path):
    """Multi-batch conversion (4 batches, batch_size=5) with correct page mapping."""
    import fitz
    from pptx import Presentation as PptxPresentation

    from src.api_client import call_llm
    from src.artifacts import ArtifactStore
    from src.logging_config import setup_logging
    from src.message_builder import build_messages
    from src.pdf_preprocessor import pdf_to_images
    from src.pptx_assembler import PPTXAssembler

    setup_logging("WARNING")
    os.chdir(tmp_path)

    total_pages = 16
    batch_size = 5
    pdf_path = str(tmp_path / "multi.pdf")
    doc = fitz.open()
    for i in range(total_pages):
        page = doc.new_page(width=720, height=405)
        page.insert_text((100, 200), f"Page {i}", fontsize=20)
    doc.save(pdf_path)
    doc.close()

    pages = pdf_to_images(pdf_path, dpi=96)
    assert len(pages) == total_pages

    store = ArtifactStore(pdf_path=pdf_path)
    store.save_page_images(pages)

    slide_xmls: dict[int, str] = {}
    batches: list[tuple[int, int]] = []
    for s in range(0, total_pages, batch_size):
        batches.append((s, min(s + batch_size, total_pages)))

    assert len(batches) == 4
    store.set_batch_count(len(batches))

    for batch_idx, (start, end) in enumerate(batches):
        batch_pages = pages[start:end]
        batch_page_map = {
            i: pages[start + i][1]["page_num"] for i in range(len(batch_pages))
        }

        messages = build_messages(batch_pages, enable_animations=False)
        store.save_messages(messages, batch_idx=batch_idx)

        stream_log = os.path.join(store.root, f"stream_batch{batch_idx}.log")
        result = call_llm(
            messages=messages,
            api_cfg=ApiConfig(api_base_url=MOCK_BASE_URL, api_key="test-key", model_name="mock-gpt-5.4"),
            stream_log_path=stream_log,
            reasoning_effort="",
            on_slide_ready=lambda pn, xml, m=batch_page_map: store.save_slide_xml(m.get(pn, pn), xml),
        )

        remapped = {batch_page_map.get(k, k): v for k, v in result.slide_xmls.items()}
        slide_xmls.update(remapped)

        store.save_api_response(result.response_data, batch_idx=batch_idx)
        store.save_stream_chunks(result.raw_chunks, batch_idx=batch_idx)
        store.save_tool_calls(result.tool_calls_raw, batch_idx=batch_idx)

    assert len(slide_xmls) == total_pages
    for i in range(total_pages):
        assert i in slide_xmls, f"Page {i} missing from slide_xmls"

    for i in range(total_pages):
        xml_path = os.path.join(store.slides_dir, f"slide_{i:03d}.xml")
        assert os.path.isfile(xml_path), f"slide_{i:03d}.xml not on disk"

    for batch_idx in range(len(batches)):
        suffix = store.batch_suffix(batch_idx)
        assert os.path.isfile(os.path.join(store.root, f"messages{suffix}.json"))
        assert os.path.isfile(os.path.join(store.root, f"api_response{suffix}.json"))
        assert os.path.isfile(os.path.join(store.root, f"stream_chunks{suffix}.jsonl"))
        assert os.path.isfile(os.path.join(store.root, f"tool_calls{suffix}.json"))

    assembler = PPTXAssembler(
        slide_width_pt=pages[0][1]["width_pt"],
        slide_height_pt=pages[0][1]["height_pt"],
    )
    assembler.assemble(slide_xmls)
    output_pptx = str(tmp_path / "multi_output.pptx")
    assembler.save(output_pptx)

    prs = PptxPresentation(output_pptx)
    assert len(prs.slides) == total_pages

    for i in range(total_pages):
        slide = prs.slides[i]
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        assert any("Test Slide Page" in t for t in texts), (
            f"Slide {i} has no content: {texts}"
        )

    shutil.rmtree("runs")
    os.remove(output_pptx)


def test_folder_input_conversion(mock_server, tmp_path):
    """Full conversion from a folder of slide images should produce a valid PPTX."""
    os.chdir(tmp_path)
    from PIL import Image

    from src.api_client import call_llm
    from src.artifacts import ArtifactStore
    from src.logging_config import setup_logging
    from src.message_builder import build_messages
    from src.pdf_preprocessor import images_from_folder, snap_slide_dimensions
    from src.pptx_assembler import PPTXAssembler
    from src.token_estimator import estimate_tokens

    setup_logging("WARNING")

    img_dir = str(tmp_path / "slides_folder")
    os.makedirs(img_dir)
    for i in range(3):
        img = Image.new("RGB", (960, 540), color=(200 + i * 20, 100, 50))
        img.save(os.path.join(img_dir, f"slide_{i:02d}.png"))

    pages = images_from_folder(img_dir)
    assert len(pages) == 3
    for idx, (_img_bytes, meta) in enumerate(pages):
        assert meta["page_num"] == idx
        assert meta["width_px"] == 960
        assert meta["height_px"] == 540
        assert "source_file" in meta

    store = ArtifactStore(pdf_path=img_dir)
    store.save_page_images(pages)

    raw_w = pages[0][1]["width_pt"]
    raw_h = pages[0][1]["height_pt"]
    snap_w, snap_h, _ratio_label = snap_slide_dimensions(raw_w, raw_h)

    messages = build_messages(pages, enable_animations=False, prompt_lang="en", provider="openai")
    store.save_messages(messages, batch_idx=0)

    token_est = estimate_tokens(messages, model="gpt-5.4")
    store.save_token_estimate(token_est)

    result = call_llm(
        messages=messages,
        api_cfg=ApiConfig(api_base_url=MOCK_BASE_URL, api_key="test-key", model_name="gpt-5.4"),
        max_tokens=4096,
        reasoning_effort="",
    )

    assert len(result.slide_xmls) == 3
    store.save_slide_xmls(result.slide_xmls)

    assembler = PPTXAssembler(slide_width_pt=snap_w, slide_height_pt=snap_h)
    assembler.assemble(result.slide_xmls)
    output_pptx = str(tmp_path / "folder_output.pptx")
    assembler.save(output_pptx)

    prs = Presentation(output_pptx)
    assert len(prs.slides) == 3
    for i in range(3):
        slide = prs.slides[i]
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        assert any("Test Slide Page" in t for t in texts)

    assert os.path.isdir(store.pages_dir)
    assert os.path.isdir(store.slides_dir)

    shutil.rmtree("runs")
    os.remove(output_pptx)


def test_folder_input_dry_run(tmp_path):
    """Dry-run from a folder of images should produce artifacts without calling the API."""
    os.chdir(tmp_path)
    from PIL import Image

    from src.dry_run import run_dry
    from src.logging_config import setup_logging

    setup_logging("WARNING")

    img_dir = str(tmp_path / "dry_slides")
    os.makedirs(img_dir)
    for i in range(2):
        img = Image.new("RGB", (960, 540), color=(100, 150 + i * 30, 200))
        img.save(os.path.join(img_dir, f"page_{i:02d}.jpg"))

    output_dir = run_dry(
        pdf_path=img_dir,
        dpi=96,
        enable_animations=False,
        model_name="gpt-5.4",
        batch_size=25,
    )

    assert output_dir.startswith("runs/")
    assert os.path.isdir(output_dir)
    assert os.path.isfile(os.path.join(output_dir, "metadata.json"))
    assert os.path.isfile(os.path.join(output_dir, "messages_0.json"))
    assert os.path.isfile(os.path.join(output_dir, "token_estimate.json"))
    assert os.path.isdir(os.path.join(output_dir, "pages"))

    page_files = os.listdir(os.path.join(output_dir, "pages"))
    assert len(page_files) == 2

    shutil.rmtree("runs")
