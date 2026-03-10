"""Unit tests for src.pptx_assembler.PPTXAssembler."""

from __future__ import annotations

from pptx import Presentation

from tests.conftest import SAMPLE_SLIDE_XML


class TestPPTXAssembler:
    SLIDE_XML = SAMPLE_SLIDE_XML.format(page_num=0).replace("Test Slide Page 0", "Hello World")

    def test_assemble_single_slide(self, tmp_path):
        from src.pptx_assembler import PPTXAssembler

        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: self.SLIDE_XML})

        output = str(tmp_path / "test.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 1
        texts = [s.text for s in prs.slides[0].shapes if hasattr(s, "text")]
        assert any("Hello World" in t for t in texts)

    def test_assemble_multiple_slides(self, tmp_path):
        from src.pptx_assembler import PPTXAssembler

        slide_xmls = {
            0: self.SLIDE_XML,
            1: self.SLIDE_XML.replace("Hello World", "Second Slide"),
        }
        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble(slide_xmls)

        output = str(tmp_path / "multi.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 2

    def test_slide_dimensions(self, tmp_path):
        import pytest

        from src.pptx_assembler import PPTXAssembler

        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: self.SLIDE_XML})

        output = str(tmp_path / "dims.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert prs.slide_width / 12700 == pytest.approx(720, abs=1)
        assert prs.slide_height / 12700 == pytest.approx(405, abs=1)

    def test_handles_hyperlinks(self, tmp_path):
        from src.pptx_assembler import PPTXAssembler

        xml_with_link = self.SLIDE_XML.replace(
            '<a:rPr lang="en-US" sz="2400"/>',
            '<a:rPr lang="en-US" sz="2400"><a:hlinkClick r:id="rId2"/></a:rPr>',
        )
        assembler = PPTXAssembler(slide_width_pt=720, slide_height_pt=405)
        assembler.assemble({0: xml_with_link})

        output = str(tmp_path / "link.pptx")
        assembler.save(output)

        prs = Presentation(output)
        assert len(prs.slides) == 1
