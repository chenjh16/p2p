"""Assemble validated slide XML into a PPTX file."""

from __future__ import annotations

import copy
import re

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Pt

from .logging_config import get_logger
from .xml_validator import validate_and_fix

logger = get_logger("assembler")

_NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

_R_NS = _NSMAP["r"]


def _qn(tag: str) -> str:
    """Convert a prefixed tag like 'p:sp' to Clark notation."""
    prefix, local = tag.split(":")
    return f"{{{_NSMAP[prefix]}}}{local}"


class PPTXAssembler:
    """Assemble GPT-generated Slide XMLs into a valid PPTX file."""

    def __init__(self, slide_width_pt: float, slide_height_pt: float):
        self.prs = Presentation()
        self.prs.slide_width = Pt(slide_width_pt)
        self.prs.slide_height = Pt(slide_height_pt)
        self._blank_layout = self.prs.slide_layouts[6]

    def assemble(self, slide_xmls: dict[int, str]) -> Presentation:  # type: ignore[valid-type]
        """Inject slide XMLs into the presentation in page order."""
        w = self.prs.slide_width or 0
        h = self.prs.slide_height or 0
        logger.info("Assembling PPTX (%.0fpt × %.0fpt)...", w / 12700, h / 12700)

        for page_num in sorted(slide_xmls.keys()):
            raw_xml = slide_xmls[page_num]
            xml_str = validate_and_fix(raw_xml, page_num)

            slide = self.prs.slides.add_slide(self._blank_layout)

            try:
                new_sld = etree.fromstring(xml_str.encode("utf-8"))
            except etree.XMLSyntaxError:
                logger.error(
                    "Slide %d: failed to parse even after validation", page_num
                )
                continue

            old_sld = slide._element

            for child in list(old_sld):
                old_sld.remove(child)

            for child in new_sld:
                old_sld.append(copy.deepcopy(child))

            for prefix, uri in new_sld.nsmap.items():
                if prefix is not None:
                    existing = old_sld.nsmap.get(prefix)
                    if existing is None or existing != uri:
                        etree.register_namespace(prefix, uri)

            n_rels = _register_relationships(slide, old_sld)

            shapes = new_sld.findall(".//" + _qn("p:sp"))
            frames = new_sld.findall(".//" + _qn("p:graphicFrame"))
            logger.info(
                "  Slide %3d: %d shapes, %d frames, %d rels",
                page_num,
                len(shapes),
                len(frames) if frames else 0,
                n_rels,
            )

        logger.info("PPTX assembled (%d slides)", len(self.prs.slides))
        return self.prs

    def save(self, output_path: str) -> None:
        """Save the assembled presentation to a PPTX file."""
        self.prs.save(output_path)
        logger.info("PPTX saved: %s", output_path)


def _register_relationships(slide, sld_element: etree._Element) -> int:  # type: ignore[name-defined]
    """Scan the slide XML for r:link/r:embed references and register them as relationships.

    Returns the number of relationships registered.
    """
    count = 0
    r_link = f"{{{_R_NS}}}link"
    r_embed = f"{{{_R_NS}}}embed"

    rid_url_map: dict[str, str] = {}

    for elem in sld_element.iter():
        for attr_name in (r_link, r_embed):
            rid = elem.get(attr_name)
            if not rid:
                continue
            if rid in rid_url_map:
                continue

            url = _extract_url_from_context(elem, rid)
            if url:
                rid_url_map[rid] = url

    for rid, url in rid_url_map.items():
        try:
            slide.part.rels.get_or_add_ext_rel(RT.HYPERLINK, url)
            count += 1
        except Exception:  # noqa: BLE001
            logger.debug("Could not register relationship %s → %s", rid, url)

    return count


def _extract_url_from_context(elem: etree._Element, rid: str) -> str:  # type: ignore[name-defined]  # pylint: disable=unused-argument
    """Try to extract a URL from an element's context (e.g. hlinkClick with an action URL)."""
    action = elem.get("action", "")
    if action.startswith("ppaction://hlinksldjump"):
        return ""
    if action.startswith("ppaction://"):
        return ""

    tooltip = elem.get("tooltip", "")
    if tooltip and (tooltip.startswith("http://") or tooltip.startswith("https://")):
        return tooltip

    tag_str = etree.tostring(elem, encoding="unicode")
    url_match = re.search(r'(https?://[^\s"<>]+)', tag_str)
    if url_match:
        return url_match.group(1)

    return ""
