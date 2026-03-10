"""Replace __LLMCLIP__ placeholders with cropped PDF raster regions."""

from __future__ import annotations

import re
from io import BytesIO

import fitz  # PyMuPDF
from pptx import Presentation

from .logging_config import get_logger

logger = get_logger("postprocessor")

_CLIP_PATTERN = re.compile(
    r"__LLMCLIP__:\[([0-9.]+),\s*([0-9.]+)\]\[([0-9.]+),\s*([0-9.]+)\]"
)


def postprocess_raster_fills(
    pptx_path: str,
    pdf_path: str,
    output_path: str,
    dpi: int = 300,
    page_indices: list[int] | None = None,
) -> None:
    """Scan PPTX for __LLMCLIP__ placeholders and fill them with cropped PDF regions.

    ``page_indices`` maps each PPTX slide position to the original PDF page
    number.  For example, if the user converted pages [0, 3, 5], then
    ``page_indices=[0, 3, 5]`` so PPTX slide 1 clips from PDF page 3.
    When *None*, a 1:1 identity mapping is assumed (slide 0 → page 0, etc.).
    """
    prs = Presentation(pptx_path)
    doc = fitz.open(pdf_path)
    fill_count = 0

    logger.info("Scanning for raster placeholders...")

    for slide_idx, slide in enumerate(prs.slides):
        if page_indices is not None:
            if slide_idx >= len(page_indices):
                logger.warning("Slide %d has no page mapping, skipping", slide_idx)
                continue
            pdf_page_num = page_indices[slide_idx]
        else:
            pdf_page_num = slide_idx

        if pdf_page_num >= doc.page_count:
            logger.warning(
                "Slide %d maps to PDF page %d which doesn't exist, skipping",
                slide_idx, pdf_page_num,
            )
            continue

        page = doc[pdf_page_num]
        shapes_to_replace: list[tuple] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            match = _CLIP_PATTERN.search(text)
            if not match:
                continue

            x1 = float(match.group(1))
            y1 = float(match.group(2))
            x2 = float(match.group(3))
            y2 = float(match.group(4))

            clip_rect = fitz.Rect(x1, y1, x2, y2)
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat, clip=clip_rect)
            img_bytes = pix.tobytes("png")

            shapes_to_replace.append((shape, img_bytes))

        for shape, img_bytes in shapes_to_replace:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height

            slide.shapes.add_picture(
                BytesIO(img_bytes), left, top, width, height
            )

            sp = shape._element
            sp.getparent().remove(sp)
            fill_count += 1

        if shapes_to_replace:
            logger.info(
                "  Slide %3d → PDF page %3d: %d placeholders filled",
                slide_idx,
                pdf_page_num,
                len(shapes_to_replace),
            )

    doc.close()
    prs.save(output_path)
    logger.info("Post-processing complete: %d raster images filled", fill_count)
